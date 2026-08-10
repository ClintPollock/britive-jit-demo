#!/usr/bin/env python3
"""
Build the 2-slide customer deck for the agentic JIT demo.

    python slides/build-slides.py            -> slides/britive-agentic-jit.pptx

Slide 1 is the scenario as it actually runs, told through the three prompts a
customer will watch being typed. Slide 2 is what it proves, as a permissions
table plus the one sentence worth remembering.

Deliberately no feature bullets: the demo is the argument, the slides just
frame it.
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1F, 0x12, 0x35)      # deep purple, headings
BODY = RGBColor(0x33, 0x33, 0x3D)
MUTED = RGBColor(0x6B, 0x6B, 0x7B)
ACCENT = RGBColor(0x6C, 0x4B, 0xB6)   # Britive-ish purple
PANEL = RGBColor(0xF4, 0xF2, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1B, 0x7F, 0x5A)
RED = RGBColor(0xB3, 0x2D, 0x2D)


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tf


def para(tf, text, *, size=14, bold=False, color=BODY, space_after=6,
         align=PP_ALIGN.LEFT, font="Segoe UI", first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def box(slide, x, y, w, h, fill, *, line=None, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.text_frame.text = ""
    return shp


def arrow(slide, x, y, w):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                               Inches(w), Inches(0.22))
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def node(slide, x, y, w, h, title, lines, *, tint=PANEL):
    box(slide, x, y, w, h, tint, line=RGBColor(0xDD, 0xD8, 0xE8))
    tf = textbox(slide, x + 0.22, y + 0.20, w - 0.44, h - 0.4)
    para(tf, title, size=14, bold=True, color=INK, space_after=4, first=True)
    for ln in lines:
        para(tf, ln, size=10.5, color=MUTED, space_after=2)


# ── Slide 1 ─────────────────────────────────────────────────────────────────
def slide_one(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Title is sized to sit on ONE line at 13.33in wide — at 28pt it wrapped and
    # shoved the whole flow down into the label below it.
    tf = textbox(s, 0.7, 0.42, 12.0, 0.9)
    para(tf, "Agentic access to customer data — no standing credentials",
         size=26, bold=True, color=INK, space_after=6, first=True)
    para(tf, "Nightly pipeline writes it. You ask Claude about it. Nothing holds a key in between.",
         size=13.5, color=MUTED)

    y = 1.55
    node(s, 0.7, y, 3.5, 1.5, "GitHub Actions",
         ["Machine identity, 07:00 daily", "OIDC → Britive → JIT credential",
          "No secrets in the repo"])
    arrow(s, 4.35, y + 0.64, 0.75)
    node(s, 5.25, y, 3.0, 1.5, "s3://…/daily/",
         ["Today's customer intake batch", "43 records · $29,322 MRR",
          "Written by an identity that", "no longer exists"], tint=WHITE)
    arrow(s, 8.45, y + 0.64, 0.75)
    node(s, 9.55, y, 3.1, 1.5, "Claude Desktop",
         ["Reads it on demand via MCP", "Credential minted per question,",
          "revoked before it answers"])

    tf = textbox(s, 0.7, 3.30, 12.0, 0.3)
    para(tf, "WHAT THE CUSTOMER WATCHES", size=10, bold=True, color=ACCENT,
         space_after=0, first=True)

    # Panel height matches its content — 2.85 left a dead band at the bottom.
    box(s, 0.7, 3.62, 11.95, 2.30, PANEL, line=RGBColor(0xE2, 0xDD, 0xEC))
    tf = textbox(s, 1.05, 3.82, 11.3, 2.0)
    para(tf, "“What did the pipeline drop today?”",
         size=15, bold=True, color=INK, space_after=2, first=True)
    para(tf, "43 records, $29,322 — different numbers every morning, because it really did run overnight.",
         size=12, color=BODY, space_after=12)
    para(tf, "“Who did that read actually run as?”",
         size=15, bold=True, color=INK, space_after=2)
    para(tf, "…/cpollock-britive-s3-readonly-role/clint.pollock@jit-zsp.com",
         size=11.5, color=GREEN, space_after=12, font="Consolas")
    para(tf, "“Now show me the same thing, as the agent itself.”",
         size=15, bold=True, color=INK, space_after=2)
    para(tf, "…/cpollock-britive-s3-readonly-role/…@iam.serviceaccount.com",
         size=11.5, color=ACCENT, space_after=0, font="Consolas")

    tf = textbox(s, 0.7, 6.18, 12.0, 0.4)
    para(tf, "Same AI. Same token. Same role. Two different identities in the audit log.",
         size=14, bold=True, color=INK, space_after=0, first=True)


# ── Slide 2 ─────────────────────────────────────────────────────────────────
def slide_two(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    tf = textbox(s, 0.7, 0.45, 12.0, 1.0)
    para(tf, "Three identities touch one file. None of them keeps a key.",
         size=28, bold=True, color=INK, space_after=6, first=True)
    para(tf, "Each one is granted just-in-time, scoped to a single job, and audited separately.",
         size=13.5, color=MUTED)

    rows = [
        ("The pipeline", "GitHub Actions, unattended",
         "Write today's batch to one bucket",
         "Cannot read anything else, cannot create infrastructure"),
        ("The agent, as itself", "Britive service identity",
         "Read that one bucket",
         "Cannot write. Cannot see any other bucket."),
        ("The agent, as you", "Same token + on-behalf-of",
         "Read exactly what you can read",
         "Loses access the moment you do"),
    ]

    # Rows were 1.28 tall on a 1.42 pitch, which pushed the third one under the
    # dark callout. Tightened so three rows + the callout fit above the fold.
    y = 1.80
    hdr = ["IDENTITY", "WHAT IT MAY DO", "WHERE IT STOPS"]
    xs, ws = [0.7, 5.05, 8.9], [4.15, 3.65, 3.75]
    for x, w, h in zip(xs, ws, hdr):
        tf = textbox(s, x, y, w, 0.3)
        para(tf, h, size=10, bold=True, color=ACCENT, space_after=0, first=True)

    y += 0.38
    ROW_H, PITCH = 1.06, 1.20
    for name, who, may, stops in rows:
        box(s, 0.7, y, 11.95, ROW_H, PANEL, line=RGBColor(0xE2, 0xDD, 0xEC))
        tf = textbox(s, xs[0] + 0.25, y, ws[0] - 0.4, ROW_H, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, name, size=14, bold=True, color=INK, space_after=3, first=True)
        para(tf, who, size=10.5, color=MUTED, space_after=0)
        tf = textbox(s, xs[1], y, ws[1] - 0.3, ROW_H, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, may, size=12, color=GREEN, space_after=0, first=True)
        tf = textbox(s, xs[2], y, ws[2] - 0.3, ROW_H, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, stops, size=12, color=RED, space_after=0, first=True)
        y += PITCH

    box(s, 0.7, 6.10, 11.95, 0.85, INK)
    tf = textbox(s, 1.05, 6.10, 11.3, 0.85, anchor=MSO_ANCHOR.MIDDLE)
    para(tf,
         "The AI reaches this data because you can — not because it holds a key of its own. "
         "Revoke you, and it loses access in the same instant.",
         size=14, bold=True, color=WHITE, space_after=0, first=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide_one(prs)
    slide_two(prs)
    out = pathlib.Path(__file__).with_name("britive-agentic-jit.pptx")
    prs.save(out)
    print(f"wrote {out}  ({out.stat().st_size:,} bytes, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
